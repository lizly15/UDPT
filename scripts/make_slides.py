#!/usr/bin/env python3
"""Tạo slide thuyết trình (.pptx) cho đồ án — thiết kế tối giản, nhất quán.

Chạy trong container: pip install python-pptx pillow; python scripts/make_slides.py
Xuất: report/slides.pptx
"""
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "report", "screenshots")
LOGO = os.path.join(ROOT, "report", "logo_hcmus.jpg")

# ---- Palette ----
NAVY = RGBColor(0x1F, 0x4E, 0x79)
NAVY2 = RGBColor(0x2C, 0x65, 0x98)
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
INK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0x65, 0x34)
FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_page = 0


def _set(run, size, color=INK, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    return tf


def para(tf, text, size, color=INK, bold=False, italic=False, bullet=False,
         align=PP_ALIGN.LEFT, space=8, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = ("•  " + text) if bullet else text
    _set(r, size, color, bold, italic)
    return p


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def footer(slide):
    global _page
    _page += 1
    tf = textbox(slide, Inches(0.5), Inches(7.02), Inches(9), Inches(0.35))
    para(tf, "Hệ thống Quản trị Kinh doanh · Ứng dụng phân tán", 9, GRAY, first=True)
    tf2 = textbox(slide, Inches(11.8), Inches(7.02), Inches(1.0), Inches(0.35))
    para(tf2, str(_page), 9, GRAY, align=PP_ALIGN.RIGHT, first=True)


def head(slide, title, kicker=None):
    """Tiêu đề slide nội dung + gạch nhấn."""
    if kicker:
        tf = textbox(slide, Inches(0.6), Inches(0.42), Inches(11), Inches(0.35))
        para(tf, kicker.upper(), 12, ACCENT, bold=True, first=True)
        top = Inches(0.78)
    else:
        top = Inches(0.5)
    tf = textbox(slide, Inches(0.6), top, Inches(12.1), Inches(0.8))
    para(tf, title, 26, NAVY, bold=True, first=True)
    rect(slide, Inches(0.62), Inches(top.inches * 914400) if False else Emu(int(top) + Inches(0.72)),
         Inches(1.1), Pt(3), ACCENT)


def new(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    if bg != WHITE:
        rect(s, 0, 0, SW, SH, bg)
    return s


def fit_image(slide, path, l, t, maxw, maxh):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    w = maxw
    h = int(w / ar)
    if h > maxh:
        h = maxh
        w = int(h * ar)
    left = l + (maxw - w) // 2
    top = t + (maxh - h) // 2
    pic = slide.shapes.add_picture(path, left, top, width=w, height=h)
    pic.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    pic.line.width = Pt(0.75)
    return pic


def sp(p):  # shortcut path
    return os.path.join(SHOTS, p)


# ================= SLIDE BUILDERS =================
def title_slide():
    s = new()
    rect(s, 0, 0, SW, Inches(0.28), NAVY)
    rect(s, 0, Inches(7.22), SW, Inches(0.28), NAVY)
    tf = textbox(s, Inches(1), Inches(0.7), Inches(11.3), Inches(0.6))
    para(tf, "TRƯỜNG ĐẠI HỌC KHOA HỌC TỰ NHIÊN – ĐHQG TP.HCM", 14, GRAY, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    para(tf, "KHOA CÔNG NGHỆ THÔNG TIN", 12, GRAY, bold=True, align=PP_ALIGN.CENTER)
    tf = textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.6),
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "HỆ THỐNG QUẢN TRỊ KINH DOANH", 40, NAVY, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    para(tf, "Xây dựng ứng dụng phân tán theo kiến trúc Microservices", 18, GRAY,
         italic=True, align=PP_ALIGN.CENTER)
    # thành viên
    tf = textbox(s, Inches(2.4), Inches(3.9), Inches(8.5), Inches(2.2))
    para(tf, "Nhóm thực hiện", 15, ACCENT, bold=True, align=PP_ALIGN.CENTER, first=True)
    for name in [
        "Vũ Huyền Thiên Lý — 21126028   (35%)",
        "Nguyễn Tuấn Kiệt — 21126027   (25%)",
        "Nguyễn Thế Thanh Long — 22127247   (25%)",
        "Nguyễn Hoàng Danh — 21126057   (15%)",
    ]:
        para(tf, name, 16, INK, align=PP_ALIGN.CENTER, space=4)
    para(tf, "GVHD: Nguyễn Trường Sơn · Phạm Minh Tú", 13, GRAY, italic=True,
         align=PP_ALIGN.CENTER, space=2)
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(6.17), Inches(6.15), height=Inches(0.9))
    notes(s, "Chào thầy và các bạn. Nhóm em trình bày đồ án Hệ thống Quản trị Kinh "
              "doanh xây dựng theo kiến trúc microservices. (~30s)")


def toc_slide():
    s = new()
    head(s, "Nội dung trình bày")
    items = [
        "1. Bài toán & Yêu cầu",
        "2. Kiến trúc & Công nghệ",
        "3. Thiết kế dữ liệu",
        "4. Giao tiếp trong hệ phân tán",
        "5. Quy trình phê duyệt (Workflow)",
        "6. Ký điện tử bất đồng bộ",
        "7. Demo giao diện & chức năng",
        "8. Xử lý bài toán phân tán",
        "9. Triển khai (Docker & Kubernetes)",
        "10. Kiểm thử & Truy vết yêu cầu",
        "11. Phân công & Kết luận",
    ]
    col1 = items[:6]
    col2 = items[6:]
    tf = textbox(s, Inches(1.0), Inches(1.9), Inches(5.6), Inches(4.6))
    for i, it in enumerate(col1):
        para(tf, it, 18, INK, bold=(i == 0 or True), space=14, first=(i == 0))
    tf = textbox(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.6))
    for i, it in enumerate(col2):
        para(tf, it, 18, INK, space=14, first=(i == 0))
    footer(s)
    notes(s, "Bài trình bày gồm 11 phần: từ bài toán, kiến trúc, thiết kế, các cơ chế "
              "phân tán, demo, đến triển khai và kiểm thử. (~40s)")


def divider(num, title, sub):
    s = new(NAVY)
    tf = textbox(s, Inches(0.9), Inches(2.2), Inches(3.2), Inches(2.5),
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, num, 96, RGBColor(0x3A, 0x6D, 0x9E), bold=True, first=True)
    rect(s, Inches(4.2), Inches(2.5), Pt(3), Inches(1.9), ACCENT)
    tf = textbox(s, Inches(4.6), Inches(2.4), Inches(8), Inches(2.2),
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 34, WHITE, bold=True, first=True)
    para(tf, sub, 16, RGBColor(0xB6, 0xC7, 0xDB), italic=True, space=2)


def content(title, bullets, kicker=None):
    s = new()
    head(s, title, kicker)
    tf = textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(4.9))
    for i, b in enumerate(bullets):
        if isinstance(b, tuple):  # (text, level)
            txt, lvl = b
            para(tf, ("    " * lvl) + txt, 17 if lvl == 0 else 15,
                 INK if lvl == 0 else GRAY, bold=(lvl == 0), bullet=(lvl == 0),
                 space=10, first=(i == 0))
        else:
            para(tf, b, 18, INK, bullet=True, space=12, first=(i == 0))
    footer(s)
    return s


def two_col(title, left_h, left, right_h, right, kicker=None):
    s = new()
    head(s, title, kicker)
    rect(s, Inches(0.8), Inches(1.9), Inches(5.75), Inches(4.7), LIGHT)
    rect(s, Inches(6.78), Inches(1.9), Inches(5.75), Inches(4.7), LIGHT)
    for x, htitle, items in [(0.8, left_h, left), (6.78, right_h, right)]:
        tf = textbox(s, Inches(x + 0.25), Inches(2.1), Inches(5.25), Inches(4.3))
        para(tf, htitle, 16, NAVY, bold=True, space=10, first=True)
        for it in items:
            para(tf, it, 15, INK, bullet=True, space=8)
    footer(s)
    return s


def image_slide(title, img, caption, bullets=None, kicker=None):
    s = new()
    head(s, title, kicker)
    if bullets:
        fit_image(s, sp(img), Inches(0.7), Inches(1.95), Inches(8.0), Inches(4.35))
        tf = textbox(s, Inches(8.95), Inches(2.0), Inches(3.7), Inches(4.4),
                     anchor=MSO_ANCHOR.MIDDLE)
        for i, b in enumerate(bullets):
            para(tf, b, 15, INK, bullet=True, space=12, first=(i == 0))
        if caption:
            tfc = textbox(s, Inches(0.7), Inches(6.35), Inches(8.0), Inches(0.35))
            para(tfc, caption, 11, GRAY, italic=True, align=PP_ALIGN.CENTER, first=True)
    else:
        fit_image(s, sp(img), Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.5))
        if caption:
            tfc = textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.4))
            para(tfc, caption, 12, GRAY, italic=True, align=PP_ALIGN.CENTER, first=True)
    footer(s)
    return s


def table_slide(title, headers, rows, kicker=None, col_w=None):
    s = new()
    head(s, title, kicker)
    nr, nc = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nr, nc, Inches(0.7), Inches(1.95),
                            Inches(11.9), Inches(0.5 + 0.42 * len(rows))).table
    if col_w:
        for i, w in enumerate(col_w):
            gt.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = gt.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = FONT
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.text = val
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            p = c.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.color.rgb = INK
            p.runs[0].font.name = FONT
    footer(s)
    return s


def closing():
    s = new(NAVY)
    tf = textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(2),
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "CẢM ƠN THẦY VÀ CÁC BẠN ĐÃ LẮNG NGHE", 30, WHITE, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    para(tf, "Mã nguồn: github.com/lizly15/UDPT", 16, RGBColor(0xB6, 0xC7, 0xDB),
         align=PP_ALIGN.CENTER, space=2)
    para(tf, "Q&A", 20, ACCENT, bold=True, align=PP_ALIGN.CENTER, space=14)


# ================= BUILD DECK =================
title_slide()
toc_slide()

# --- Phần 1 ---
divider("01", "Bài toán & Yêu cầu", "Bối cảnh, mục tiêu và phạm vi hệ thống")
content("Bối cảnh doanh nghiệp", [
    "Công ty ABC: dịch vụ logistics — khai thác cảng, vận chuyển, lưu kho, bốc xếp",
    "Mỗi tháng phát sinh nhiều hợp đồng, phụ lục, bảng giá, sản lượng, bảng thanh toán",
    "Nhiều phòng ban tham gia: Kinh doanh, Khai thác, Kế toán, Pháp chế, Giám đốc",
], kicker="Phần 1 · Bài toán")
content("Vấn đề hiện trạng", [
    "Xử lý thủ công qua email / Excel / giấy tờ → thông tin phân tán",
    "Khó kiểm soát trạng thái hồ sơ, khó biết ai đang giữ hồ sơ",
    "Khó truy vết lịch sử thay đổi; dễ sai lệch khi giá/sản lượng thay đổi",
    "→ Cần một hệ thống tập trung quản lý trọn vòng đời hồ sơ kinh doanh",
], kicker="Phần 1 · Bài toán")
two_col("Yêu cầu chức năng (mục 4)",
        "Nghiệp vụ lõi", [
            "Khách hàng & danh mục dịch vụ",
            "Hợp đồng & phụ lục",
            "Bảng giá (nhiều phiên bản)",
            "Sản lượng thực hiện",
            "Bảng thanh toán",
        ],
        "Xuyên suốt", [
            "Quy trình phê duyệt cấu hình được",
            "Ký điện tử",
            "Thông báo bất đồng bộ",
            "Nhật ký & truy vết",
            "Phân quyền theo vai trò",
        ], kicker="Phần 1 · Yêu cầu")
content("Yêu cầu kỹ thuật (mục 6)", [
    "Kiến trúc Microservices: ≥ 4 service nghiệp vụ + API Gateway",
    "Backend FastAPI, có OpenAPI/Swagger; mỗi service có DB/schema riêng",
    "PostgreSQL (dữ liệu) · Redis (cache/rate-limit) · Kafka (bất đồng bộ)",
    "Docker Compose (dev) + Kubernetes (minikube)",
    "Logging, error handling, validation, phân quyền bằng JWT",
], kicker="Phần 1 · Yêu cầu")

# --- Phần 2 ---
divider("02", "Kiến trúc & Công nghệ", "Tổng thể hệ thống và stack sử dụng")
image_slide("Kiến trúc tổng thể", "diagram-arch.png",
            "Gateway điều phối & xác thực; service gọi REST cho lệnh, phát/nhận sự kiện qua Kafka",
            kicker="Phần 2 · Kiến trúc")
table_slide("9 service chạy độc lập", ["Service", "Vai trò"], [
    ["api-gateway", "Reverse proxy, verify JWT, rate-limit, idempotency"],
    ["identity", "User, role, JWT (access/refresh), blacklist"],
    ["customer", "Khách hàng + danh mục dịch vụ"],
    ["contract", "Hợp đồng + phụ lục (CTR)"],
    ["pricing", "Bảng giá + version (PRC)"],
    ["billing", "Sản lượng + bảng thanh toán (PAY)"],
    ["workflow", "Engine phê duyệt data-driven + điều phối ký (APR)"],
    ["notification", "Thông báo + nhật ký truy vết"],
    ["mock-esign", "Nhà cung cấp ký điện tử giả lập (async)"],
], kicker="Phần 2 · Kiến trúc", col_w=[2.6, 9.3])
two_col("Công nghệ sử dụng",
        "Backend & dữ liệu", [
            "FastAPI + SQLAlchemy 2 + Pydantic",
            "PostgreSQL 16 (DB-per-service)",
            "Redis · Apache Kafka (KRaft)",
            "JWT + RBAC 7 vai trò",
        ],
        "Frontend & hạ tầng", [
            "React + Vite + TypeScript + Tailwind",
            "Docker Compose (14 container)",
            "Kubernetes (minikube)",
            "OpenAPI/Swagger tự sinh",
        ], kicker="Phần 2 · Công nghệ")
content("Quyết định thiết kế chính", [
    "Microservices theo bounded context → mỗi service một trách nhiệm rõ ràng",
    "Database-per-service → cô lập dữ liệu, không FK xuyên service",
    "Bất đồng bộ qua Kafka + Outbox → chống mất sự kiện",
    "Workflow cấu hình bằng dữ liệu → không hard-code quy trình duyệt",
    "Zero-trust nội bộ: mỗi service vẫn tự verify JWT dù đã qua Gateway",
], kicker="Phần 2 · Thiết kế")

# --- Phần 3 ---
divider("03", "Thiết kế dữ liệu", "Mô hình dữ liệu và Database-per-Service")
content("Database-per-Service", [
    "1 cụm PostgreSQL, mỗi service 1 database riêng (authdb, contractdb, …)",
    "Không FK xuyên service — liên kết bằng mã nghiệp vụ (customer_code, …)",
    "Ưu điểm: cô lập lỗi, độc lập triển khai, dễ mở rộng từng service",
    "Điểm nhấn (PAY-03): bảng thanh toán COPY cứng đơn giá tại thời điểm tính",
], kicker="Phần 3 · Dữ liệu")
image_slide("Sơ đồ ERD", "erd.png",
            "Các thực thể chính và quan hệ (theo từng service)",
            kicker="Phần 3 · Dữ liệu")

# --- Phần 4 ---
divider("04", "Giao tiếp trong hệ phân tán", "Đồng bộ (REST) và bất đồng bộ (Kafka + Outbox)")
two_col("Hai kiểu giao tiếp",
        "Đồng bộ — REST", [
            "Dùng cho lệnh cần kết quả ngay",
            "Gateway → service; service → service",
            "VD: submit hồ sơ, billing lấy giá/HĐ",
        ],
        "Bất đồng bộ — Kafka", [
            "Dùng cho lan truyền sự kiện",
            "Thông báo, audit, kết quả duyệt/ký",
            "Rã kết nối (loose coupling), chịu lỗi tốt",
        ], kicker="Phần 4 · Giao tiếp")
content("Outbox Pattern — chống mất sự kiện", [
    "Vấn đề: ghi DB thành công nhưng phát Kafka lỗi → mất sự kiện",
    "Giải pháp: ghi thay đổi + bản ghi outbox trong CÙNG một transaction",
    "Một relay đọc outbox → phát Kafka → đánh dấu đã gửi",
    "Consumer idempotent → an toàn khi phát lại",
], kicker="Phần 4 · Giao tiếp")
image_slide("Sự kiện thật trên Kafka", "03-kafka-messages.png",
            None, bullets=[
                "5 topic: contract / pricing / billing / workflow / esign",
                "Mỗi message có event_id, occurred_at",
                "Phát qua Outbox, consumer xử lý idempotent",
            ], kicker="Phần 4 · Giao tiếp")

# --- Phần 5 ---
divider("05", "Quy trình phê duyệt", "Workflow engine cấu hình bằng dữ liệu")
content("Workflow data-driven (không hard-code)", [
    "Quy trình lưu trong bảng: định nghĩa theo loại hồ sơ + các bước + người duyệt",
    "Engine đọc DB để quyết định bước kế tiếp — KHÔNG viết if/else theo loại hồ sơ",
    "Mỗi loại hồ sơ (Hợp đồng / Bảng giá / Thanh toán) có quy trình riêng",
    "Thêm/sửa quy trình = sửa dữ liệu, không cần sửa code",
], kicker="Phần 5 · Workflow")
image_slide("Cấu hình quy trình duyệt", "17-admin.png",
            None, bullets=[
                "Quản trị người dùng + vai trò",
                "Xem cấu hình từng quy trình",
                "Hợp đồng: 4 cấp; Thanh toán: 2 cấp",
            ], kicker="Phần 5 · Workflow")
image_slide("Trình tự duyệt hợp đồng", "diagram-seq-contract.png",
            "Submit → duyệt các cấp → phát DocApproved (Kafka) → hồ sơ chuyển Approved",
            kicker="Phần 5 · Workflow")
content("Các quy tắc phê duyệt (APR)", [
    "APR-01: chỉ đúng người được giao ở bước hiện tại (không chỉ kiểm role)",
    "APR-02: không duyệt nhảy bước / duyệt lại bước đã xong",
    "APR-03: reject / yêu cầu sửa bắt buộc nhập lý do",
    "APR-05: bước cuối duyệt → phát sự kiện để xử lý tiếp (vd gửi ký)",
], kicker="Phần 5 · Workflow")

# --- Phần 6 ---
divider("06", "Ký điện tử", "Tích hợp bất đồng bộ với callback")
image_slide("Luồng ký điện tử bất đồng bộ", "diagram-seq-payment.png",
            "Bảng thanh toán duyệt xong → workflow tự gửi ký → mock-esign callback → Issued",
            kicker="Phần 6 · Ký điện tử")
content("Quản lý phiên ký & xử lý lỗi", [
    "Trạng thái phiên ký TÁCH BIỆT trạng thái duyệt (pending/signing/signed/failed)",
    "Chỉ gửi ký sau khi hồ sơ được duyệt nội bộ (PAY-06)",
    "Ký thất bại → phản ánh rõ + cho phép GỬI KÝ LẠI (PAY-07, SC-06)",
    "Nếu dịch vụ ký lỗi tạm thời → có retry, không hỏng dữ liệu đã duyệt",
], kicker="Phần 6 · Ký điện tử")

# --- Phần 7 ---
divider("07", "Demo giao diện & chức năng", "Minh chứng các chức năng nghiệp vụ")
image_slide("Đăng nhập & Tổng quan", "06-dashboard.png",
            "Xác thực JWT; dashboard tổng hợp số liệu theo vai trò",
            kicker="Phần 7 · Demo")
image_slide("Quản lý khách hàng", "07-customers.png",
            "Chức năng 4.1 — CRUD khách hàng, danh mục dịch vụ",
            kicker="Phần 7 · Demo")
image_slide("Hợp đồng & tiến trình duyệt", "09-contract-detail.png",
            "Chức năng 4.2/4.3 — chi tiết hợp đồng + tiến trình duyệt 4 cấp",
            kicker="Phần 7 · Demo")
image_slide("Bảng giá nhiều phiên bản", "10-pricing.png",
            "Chức năng 4.4 — version Effective/Superseded, tra giá theo ngày",
            kicker="Phần 7 · Demo")
image_slide("Sản lượng & khóa kỳ", "11-volumes.png",
            "Chức năng 4.5 — ghi nhận sản lượng, khóa kỳ trước khi tính phí",
            kicker="Phần 7 · Demo")
image_slide("Bảng thanh toán (khớp Phụ lục A.8)", "12-payment-detail.png",
            "Chức năng 4.6/4.8 — tổng 50.500.000đ, trạng thái Issued, ký điện tử signed",
            kicker="Phần 7 · Demo")
image_slide("Hộp thư phê duyệt", "14-approval-inbox.png",
            "Chức năng 4.7 — chỉ đúng người được giao mới thấy & xử lý tác vụ",
            kicker="Phần 7 · Demo")
image_slide("Thông báo & Nhật ký", "16-audit.png",
            "Chức năng 4.9/4.10 — thông báo bất đồng bộ + truy vết đầy đủ theo hồ sơ",
            kicker="Phần 7 · Demo")

# --- Phần 8 ---
divider("08", "Xử lý bài toán phân tán", "Các điểm nhấn kỹ thuật")
table_slide("7 vấn đề phân tán & giải pháp", ["Vấn đề", "Giải pháp"], [
    ["Quy trình không hard-code", "Engine đọc định nghĩa từ DB"],
    ["Double submit (SC-09)", "Idempotency-Key (Redis) + instance duy nhất"],
    ["Race condition duyệt (SC-05)", "Optimistic locking (version) → 409"],
    ["Mất sự kiện", "Outbox Pattern + consumer idempotent"],
    ["Phân quyền theo ngữ cảnh (SC-08)", "Kiểm tra đúng assignee, không chỉ role"],
    ["Dữ liệu lịch sử (SC-04/10)", "Chọn giá theo ngày + copy đơn giá cứng"],
    ["Service phụ lỗi (SC-07)", "Async + Outbox, nghiệp vụ chính không hỏng"],
], kicker="Phần 8 · Phân tán", col_w=[5.2, 6.7])
image_slide("Chặn sai quy tắc nghiệp vụ", "13-sc01-error.png",
            None, bullets=[
                "CTR-02/SC-01: thiếu tài liệu → chặn submit",
                "Rule enforce ở tầng service",
                "Trả mã lỗi rõ ràng, chuẩn hóa",
            ], kicker="Phần 8 · Phân tán")

# --- Phần 9 ---
divider("09", "Triển khai", "Docker Compose và Kubernetes")
image_slide("Docker Compose (môi trường dev)", "01-docker-ps.png",
            None, bullets=[
                "14 container chạy độc lập",
                "1 lệnh: make up",
                "Postgres · Redis · Kafka · 9 service",
            ], kicker="Phần 9 · Triển khai")
image_slide("Kubernetes (minikube)", "20-k8s-pods.png",
            None, bullets=[
                "Manifests: namespace, config, hạ tầng, 9 service",
                "Toàn bộ pod Running",
                "Deploy được trên minikube",
            ], kicker="Phần 9 · Triển khai")

# --- Phần 10 ---
divider("10", "Kiểm thử & Truy vết", "Chất lượng và bao phủ yêu cầu")
image_slide("Kiểm thử kịch bản (SC-01→SC-10)", "18-smoke-test.png",
            None, bullets=[
                "Smoke test end-to-end",
                "14/14 assertion PASS",
                "Bao phủ toàn bộ Phụ lục A.12",
            ], kicker="Phần 10 · Kiểm thử")
image_slide("Unit test (pytest)", "19-unit-test.png",
            None, bullets=[
                "State machine & business rules",
                "31/31 test PASS",
                "4 service: contract, pricing, workflow, billing",
            ], kicker="Phần 10 · Kiểm thử")
content("Bao phủ yêu cầu (truy vết)", [
    "Chức năng 4.1–4.10: đầy đủ, có minh chứng giao diện",
    "Business rules CTR / PRC / PAY / APR: enforce ở tầng service",
    "Kịch bản SC-01→SC-10: 14/14 PASS (tự động)",
    "Yêu cầu kỹ thuật mục 6: Microservices, Gateway, PostgreSQL, Redis, Kafka, "
    "Docker, Kubernetes, JWT — đủ",
], kicker="Phần 10 · Truy vết")

# --- Phần 11 ---
divider("11", "Kết luận", "Tổng kết và hướng phát triển")
content("Kết luận", [
    "Hiện thực đầy đủ nghiệp vụ quản trị kinh doanh trên kiến trúc microservices",
    "Giải quyết trọn các bài toán phân tán: Outbox, idempotency, optimistic lock, "
    "phân quyền ngữ cảnh, ký điện tử async",
    "Triển khai thành công trên Docker Compose và Kubernetes",
    "Kiểm chứng bằng bộ test tự động: 14/14 smoke + 31/31 unit",
], kicker="Phần 11 · Kết luận")
content("Hướng phát triển", [
    "Thay Outbox relay bằng CDC (Debezium)",
    "Thêm distributed tracing (OpenTelemetry) để quan sát xuyên service",
    "Auto-scaling (HPA) và lưu trữ bền (PVC/StatefulSet) trên Kubernetes",
    "Tích hợp dịch vụ ký điện tử thật thay cho mock",
], kicker="Phần 11 · Kết luận")

closing()

out = os.path.join(ROOT, "report", "slides.pptx")
prs.save(out)
print(f"✅ Đã tạo {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slide")
